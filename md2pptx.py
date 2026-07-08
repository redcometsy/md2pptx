"""Markdown to PPTX converter using python-pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re
import sys


def parse_markdown(md_text):
    """Parse markdown into slides. Each ## heading starts a new slide."""
    lines = md_text.strip().split('\n')
    slides = []
    current = None

    for line in lines:
        # H1 = title slide
        if line.startswith('# ') and not line.startswith('## '):
            slides.append({'title': line[2:].strip(), 'type': 'title', 'content': []})
        # H2 = new slide
        elif line.startswith('## '):
            if current:
                slides.append(current)
            current = {'title': line[3:].strip(), 'type': 'slide', 'content': []}
        # H3 = sub-item in current slide
        elif line.startswith('### '):
            if current is None:
                current = {'title': line[4:].strip(), 'type': 'slide', 'content': []}
            else:
                current['content'].append({'type': 'h3', 'text': line[4:].strip()})
        # Unordered list
        elif line.strip().startswith('- ') or line.strip().startswith('* '):
            text = re.sub(r'^[\s]*[-*]\s+', '', line).strip()
            if current is None:
                current = {'title': text, 'type': 'slide', 'content': []}
            else:
                current['content'].append({'type': 'bullet', 'text': text})
        # Code block
        elif line.strip().startswith('```'):
            continue  # skip code fences in simple mode
        # Regular text
        elif line.strip():
            if current is None:
                current = {'title': line.strip(), 'type': 'slide', 'content': []}
            else:
                current['content'].append({'type': 'text', 'text': line.strip()})

    if current:
        slides.append(current)

    return slides


def create_pptx(slides, output_path):
    """Create PPTX from parsed slides."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    for i, slide_data in enumerate(slides):
        if slide_data['type'] == 'title':
            # Title slide
            slide_layout = prs.slide_layouts[6]  # blank
            slide = prs.slides.add_slide(slide_layout)

            # Background
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0x2B, 0x2B, 0x2B)

            # Title
            left = Inches(1)
            top = Inches(2.5)
            width = Inches(11.333)
            height = Inches(2)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data['title']
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.alignment = PP_ALIGN.CENTER

        else:
            # Content slide
            slide_layout = prs.slide_layouts[6]  # blank
            slide = prs.slides.add_slide(slide_layout)

            # Background
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0x2B, 0x2B, 0x2B)

            # Slide number bar
            left = Inches(0)
            top = Inches(0)
            width = Inches(13.333)
            height = Inches(0.06)
            shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0x00, 0xBC, 0x7D)
            shape.line.fill.background()

            # Title
            left = Inches(0.8)
            top = Inches(0.4)
            width = Inches(11.733)
            height = Inches(1)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data['title']
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            # Content
            left = Inches(0.8)
            top = Inches(1.6)
            width = Inches(11.733)
            height = Inches(5.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True

            for j, item in enumerate(slide_data['content']):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                if item['type'] == 'bullet':
                    p.text = f"  {item['text']}"
                    p.font.size = Pt(22)
                    p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
                    p.space_after = Pt(8)
                    # Add bullet character
                    p.level = 0
                elif item['type'] == 'h3':
                    p.text = item['text']
                    p.font.size = Pt(26)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0x00, 0xBC, 0x7D)
                    p.space_after = Pt(12)
                    p.space_before = Pt(16)
                else:
                    p.text = item['text']
                    p.font.size = Pt(22)
                    p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
                    p.space_after = Pt(6)

        # Page number
        if slide_data['type'] != 'title':
            left = Inches(12.2)
            top = Inches(7.0)
            width = Inches(0.8)
            height = Inches(0.4)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = str(i)
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
            p.alignment = PP_ALIGN.RIGHT

    prs.save(output_path)
    print(f"Saved: {output_path} ({len(slides)} slides)")


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Usage: python md2pptx.py <input.md> [output.pptx]")
        sys.exit(1)

    input_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else input_path.rsplit('.', 1)[0] + '.pptx'

    with open(input_path, 'r', encoding='utf-8') as f:
        md_text = f.read()

    slides = parse_markdown(md_text)
    create_pptx(slides, output_path)
