"""PPTX to Markdown converter — reverse of md2pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt
import sys


def is_page_number(shape):
    """Check if shape is a page number (bottom-right, small single-line text)."""
    top = shape.top or 0
    left = shape.left or 0
    if top < Inches(6.5) or left < Inches(11):
        return False
    if not shape.has_text_frame:
        return False
    text = shape.text_frame.text.strip()
    return text.isdigit() and len(text) <= 2


def pptx_to_markdown(input_path):
    """Convert PPTX to markdown string."""
    prs = Presentation(input_path)
    lines = []

    for i, slide in enumerate(prs.slides):
        shapes_by_pos = []

        for shape in slide.shapes:
            # Skip page numbers (bottom-right corner)
            if is_page_number(shape):
                continue
            shapes_by_pos.append((shape.top or 0, shape.left or 0, shape))

        shapes_by_pos.sort(key=lambda x: (x[0], x[1]))

        if not shapes_by_pos:
            continue

        # Separate title, content, skip empty/decorative shapes
        title_text = ""
        content_items = []
        for _, _, shape in shapes_by_pos:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if not title_text:
                title_text = text
            else:
                # Multiple paragraphs in content textbox
                for para in shape.text_frame.paragraphs:
                    pt = para.text.strip()
                    if pt:
                        content_items.append(pt)

        if i == 0 and len(shapes_by_pos) <= 2:
            lines.append(f"# {title_text}")
        else:
            lines.append(f"## {title_text}")

        lines.append("")

        for item in content_items:
            lines.append(f"- {item}")

        lines.append("")

    return "\n".join(lines)


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Usage: python pptx2md.py <input.pptx> [output.md]")
        sys.exit(1)

    input_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else input_path.rsplit('.', 1)[0] + '.md'

    md_text = pptx_to_markdown(input_path)

    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(md_text)

    slide_count = len(Presentation(input_path).slides)
    print(f"Saved: {output_path} ({len(md_text.splitlines())} lines, {slide_count} slides)")
